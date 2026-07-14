"""
Genera la presentación de ALBA para el Concurso Datos al Ecosistema 2026.
Crea:
  - recursos/presentacion.pptx
  - recursos/presentacion.pdf
  - recursos/portada.png
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
from PIL import Image, ImageDraw, ImageFont
import sys

ROOT = Path(__file__).resolve().parent.parent
RESOURCES = ROOT / "recursos"
RESOURCES.mkdir(exist_ok=True)

# Colores ALBA
DARK_BG = RGBColor(0x05, 0x08, 0x13)
GOLD = RGBColor(0xD4, 0xAF, 0x37)
SLATE = RGBColor(0x94, 0xA3, 0xB8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x34, 0xD3, 0x99)
ROSE = RGBColor(0xF4, 0x71, 0x71)


def add_background(slide, color=DARK_BG):
    """Añade un rectángulo de fondo que cubre toda la diapositiva."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_title_box(slide, text, left, top, width, height, font_size=32, bold=True, color=WHITE):
    tf = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height)).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.LEFT
    return tf


def add_bullet_box(slide, items, left, top, width, height, font_size=16, color=SLATE, bullet_color=GOLD):
    tf = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height)).text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(10)
        p.level = 0
    return tf


def add_footer(slide, text="ALBA — Datos al Ecosistema 2026"):
    tf = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12), Inches(0.4)).text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = GOLD
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.RIGHT


# ---------------------------------------------------------------------------
# 1. Crear presentación
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]  # blank

# ---------------------------------------------------------------------------
# Diapositiva 1: Portada
# ---------------------------------------------------------------------------
slide = prs.slides.add_slide(blank_layout)
add_background(slide)
# Línea dorada decorativa
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.4), Inches(1.5), Inches(0.05))
line.fill.solid()
line.fill.fore_color.rgb = GOLD
line.line.fill.background()

add_title_box(slide, "ALBA", 0.5, 1.2, 12, 1.2, font_size=72, bold=True, color=GOLD)
add_title_box(slide, "Analítica Laboral Basada en IA", 0.5, 2.6, 12, 0.8, font_size=28, bold=False, color=WHITE)
add_title_box(slide, "Reto 5: Economía y Empleo — Concurso Datos al Ecosistema 2026", 0.5, 3.4, 12, 0.6, font_size=18, bold=False, color=SLATE)
add_title_box(slide, "Nivel: Avanzado", 0.5, 4.2, 12, 0.6, font_size=16, bold=False, color=GOLD)
add_title_box(slide, "Manuel Francisco Machado\nLíder / Desarrollador", 0.5, 5.2, 12, 1, font_size=16, bold=False, color=SLATE)
add_footer(slide, "Plataforma de inteligencia laboral para Colombia")

# ---------------------------------------------------------------------------
# Diapositiva 2: Problema
# ---------------------------------------------------------------------------
slide = prs.slides.add_slide(blank_layout)
add_background(slide)
add_title_box(slide, "El problema", 0.5, 0.5, 12, 0.8, font_size=40, color=WHITE)
add_bullet_box(slide, [
    "24 millones de ocupados, pero 58% en informalidad (GEIH 2025).",
    "428.000 programas académicos matriculados, muchos sin correspondencia con la demanda laboral.",
    "Datos dispersos entre DANE, MEN, SENA, MinTrabajo, DNP y Confecámaras.",
    "Ciudadanos no saben qué estudiar; emprendedores no saben qué negocio crear;",
    "empresas no encuentran talento; el gobierno no tiene visibilidad integral."
], 0.5, 1.5, 12, 5, font_size=20, color=SLATE)
add_footer(slide)

# ---------------------------------------------------------------------------
# Diapositiva 3: Datos abiertos
# ---------------------------------------------------------------------------
slide = prs.slides.add_slide(blank_layout)
add_background(slide)
add_title_box(slide, "Datos abiertos utilizados", 0.5, 0.5, 12, 0.8, font_size=40, color=WHITE)
add_bullet_box(slide, [
    "12 fuentes oficiales · 44 tablas · ~744.000 filas en Supabase",
    "DANE GEIH: empleo, salarios, informalidad por departamento y ocupación",
    "MinTrabajo PILA: cotizantes formales por sector CIIU",
    "Confecámaras RUES: empresas nuevas por sector y territorio",
    "MEN SNIES + OLE: oferta educativa e ingresos reales de egresados",
    "SENA SPE/APE: demanda laboral y cursos de formación",
    "ESCO, O*NET y World Bank: habilidades, ocupaciones e indicadores macro"
], 0.5, 1.5, 8.5, 5, font_size=18, color=SLATE)
# Cuadro resumen
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.2), Inches(1.5), Inches(3.6), Inches(4.8))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(0x0A, 0x0F, 0x1F)
box.line.color.rgb = GOLD
box.line.width = Pt(1.5)
add_title_box(slide, "11\nFuentes", 9.4, 1.8, 1.5, 1, font_size=28, color=GOLD)
add_title_box(slide, "44\nTablas", 11.0, 1.8, 1.5, 1, font_size=28, color=GOLD)
add_title_box(slide, "744K\nFilas", 9.4, 3.2, 1.5, 1, font_size=28, color=GOLD)
add_title_box(slide, "33\nDeptos", 11.0, 3.2, 1.5, 1, font_size=28, color=GOLD)
add_title_box(slide, "406\nSalarios reales", 9.4, 4.6, 3.0, 0.8, font_size=20, color=WHITE)
add_footer(slide)

# ---------------------------------------------------------------------------
# Diapositiva 4: Solución e IA
# ---------------------------------------------------------------------------
slide = prs.slides.add_slide(blank_layout)
add_background(slide)
add_title_box(slide, "Solución e Inteligencia Artificial", 0.5, 0.5, 12, 0.8, font_size=40, color=WHITE)
add_bullet_box(slide, [
    "Observatorio Inteligente: panorama real del mercado laboral colombiano.",
    "Predicción IA: proyecciones a 5 y 10 años con Chronos T5 Small.",
    "Match Inteligente: une CV, vacante, ESCO, salarios de egresados y LLM.",
    "Emprende IA: índice de oportunidad por municipio con RUES + EMICRON.",
    "Coach IA: mejora CV y simula entrevistas con Gemini 2.5 Flash-Lite.",
    "Simulación: ¿qué pasa si cambio de carrera o sector?"
], 0.5, 1.5, 8.5, 5, font_size=18, color=SLATE)
# Modelos
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.2), Inches(1.5), Inches(3.6), Inches(4.8))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(0x0A, 0x0F, 0x1F)
box.line.color.rgb = GOLD
box.line.width = Pt(1.5)
add_title_box(slide, "Modelos IA", 9.4, 1.7, 3.2, 0.6, font_size=20, color=GOLD)
add_title_box(slide, "Gemini 2.5\nFlash-Lite", 9.4, 2.5, 3.2, 0.8, font_size=16, color=WHITE)
add_title_box(slide, "Gemini Live", 9.4, 3.4, 3.2, 0.6, font_size=16, color=WHITE)
add_title_box(slide, "Chronos T5\nSmall", 9.4, 4.2, 3.2, 0.8, font_size=16, color=WHITE)
add_title_box(slide, "Gemma\nEmbeddings 300", 9.4, 5.0, 3.2, 0.8, font_size=16, color=WHITE)
add_footer(slide)

# ---------------------------------------------------------------------------
# Diapositiva 5: Arquitectura
# ---------------------------------------------------------------------------
slide = prs.slides.add_slide(blank_layout)
add_background(slide)
add_title_box(slide, "Arquitectura tecnológica", 0.5, 0.5, 12, 0.8, font_size=40, color=WHITE)
# Diagrama simplificado en capas
capas = [
    ("Frontend", "React 18 + TypeScript + Vite + Tailwind CSS", 1.2),
    ("Backend", "FastAPI + Uvicorn · 6 routers · 40+ endpoints", 2.4),
    ("Datos + IA", "Supabase PostgreSQL + pgvector · Gemini · Chronos T5", 3.6),
    ("Fuentes", "DANE · MEN · SENA · MinTrabajo · DNP · RUES · World Bank", 4.8),
]
for i, (titulo, desc, top) in enumerate(capas):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(top), Inches(12.3), Inches(0.9))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x0A, 0x0F, 0x1F)
    box.line.color.rgb = GOLD if i in (0, 3) else SLATE
    box.line.width = Pt(1)
    add_title_box(slide, titulo, 0.8, top + 0.15, 2.5, 0.6, font_size=18, bold=True, color=GOLD)
    add_title_box(slide, desc, 3.4, top + 0.2, 9, 0.5, font_size=15, bold=False, color=WHITE)
add_footer(slide)

# ---------------------------------------------------------------------------
# Diapositiva 6: Resultados / Demo
# ---------------------------------------------------------------------------
slide = prs.slides.add_slide(blank_layout)
add_background(slide)
add_title_box(slide, "Resultados y demo", 0.5, 0.5, 12, 0.8, font_size=40, color=WHITE)
add_bullet_box(slide, [
    "Integración de 11 fuentes en 44 tablas y ~744.000 filas.",
    "Predicciones a 5 y 10 años para sectores, profesiones, habilidades y salarios.",
    "Salarios reales de 406 ocupaciones del DANE GEIH (2026-04).",
    "Matching híbrido: 50% habilidades ESCO + 50% análisis del LLM.",
    "Mapa territorial con empleo por departamento y sector.",
    "Botón 'Analizar con IA' en cada gráfica del Observatorio y Predicción."
], 0.5, 1.5, 12, 4.5, font_size=20, color=SLATE)
# Métricas
metrics = [
    ("5", "Módulos"),
    ("40+", "Endpoints"),
    ("406", "Ocupaciones con salario real"),
    ("21", "Profesiones proyectadas"),
]
for i, (valor, label) in enumerate(metrics):
    x = 0.5 + i * 3.1
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(5.8), Inches(2.8), Inches(1.2))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x0A, 0x0F, 0x1F)
    box.line.color.rgb = GOLD
    add_title_box(slide, valor, x + 0.1, 6.0, 2.6, 0.5, font_size=26, color=GOLD)
    add_title_box(slide, label, x + 0.1, 6.45, 2.6, 0.4, font_size=12, color=SLATE)
add_footer(slide)

# ---------------------------------------------------------------------------
# Diapositiva 7: Impacto
# ---------------------------------------------------------------------------
slide = prs.slides.add_slide(blank_layout)
add_background(slide)
add_title_box(slide, "Impacto esperado", 0.5, 0.5, 12, 0.8, font_size=40, color=WHITE)
add_bullet_box(slide, [
    "Ciudadanos: toman decisiones informadas sobre qué estudiar y dónde buscar empleo.",
    "Emprendedores: identifican negocios con mayor potencial en su municipio.",
    "Universidades: evalúan la alineación de sus programas con el mercado laboral.",
    "Empresas: encuentran talento y entienden la demanda de su sector.",
    "Gobierno: diseña políticas de empleo y formación con visibilidad integral.",
    "Escalabilidad: arquitectura lista para desplegar en Vercel + Railway + Supabase."
], 0.5, 1.5, 12, 5, font_size=20, color=SLATE)
add_footer(slide)

# ---------------------------------------------------------------------------
# Diapositiva 8: Repositorio y cierre
# ---------------------------------------------------------------------------
slide = prs.slides.add_slide(blank_layout)
add_background(slide)
add_title_box(slide, "Repositorio y valor diferencial", 0.5, 0.5, 12, 0.8, font_size=40, color=WHITE)
add_bullet_box(slide, [
    "Repositorio GitHub con código fuente organizado.",
    "Docs: arquitectura, diccionario de datos, metodología, fuentes y conclusiones.",
    "Demo en vivo: app web con backend FastAPI y base de datos Supabase.",
    "Valor diferencial: datos reales del DANE + IA explicable + análisis por territorio.",
    "Datos abiertos + IA = decisiones laborales más inteligentes para Colombia."
], 0.5, 1.5, 12, 4.5, font_size=20, color=SLATE)
add_title_box(slide, "¡Gracias!", 0.5, 6.0, 12, 0.8, font_size=48, bold=True, color=GOLD)
add_footer(slide)

# Guardar PPTX
pptx_path = RESOURCES / "presentacion.pptx"
prs.save(str(pptx_path))
print(f"Guardado: {pptx_path}")

# ---------------------------------------------------------------------------
# 2. Exportar a PDF con PowerPoint via comtypes
# ---------------------------------------------------------------------------
pdf_path = RESOURCES / "presentacion.pdf"
try:
    import comtypes.client
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1
    deck = powerpoint.Presentations.Open(str(pptx_path), WithWindow=False)
    deck.SaveAs(str(pdf_path), 32)  # ppSaveAsPDF = 32
    deck.Close()
    powerpoint.Quit()
    print(f"Guardado PDF: {pdf_path}")
except Exception as e:
    print(f"No se pudo exportar a PDF automáticamente: {e}")
    print("Abre presentacion.pptx en PowerPoint y exporta manualmente a PDF.")

# ---------------------------------------------------------------------------
# 3. Generar portada.png
# ---------------------------------------------------------------------------
portada_path = RESOURCES / "portada.png"
img = Image.new("RGB", (1920, 1080), color=(0x05, 0x08, 0x13))
draw = ImageDraw.Draw(img)
# Intentar cargar fuentes
font_dir = Path("C:/Windows/Fonts")
try:
    font_title = ImageFont.truetype(str(font_dir / "Calibri.ttf"), 120)
    font_sub = ImageFont.truetype(str(font_dir / "Calibri.ttf"), 48)
    font_small = ImageFont.truetype(str(font_dir / "Calibri.ttf"), 32)
    font_tiny = ImageFont.truetype(str(font_dir / "Calibri.ttf"), 24)
except Exception:
    font_title = ImageFont.load_default()
    font_sub = font_small = font_tiny = font_title

# Línea dorada
draw.rectangle([80, 270, 300, 280], fill=(0xD4, 0xAF, 0x37))
# Textos
draw.text((80, 130), "ALBA", font=font_title, fill=(0xD4, 0xAF, 0x37))
draw.text((80, 310), "Analítica Laboral Basada en IA", font=font_sub, fill=(0xFF, 0xFF, 0xFF))
draw.text((80, 400), "Reto 5: Economía y Empleo — Concurso Datos al Ecosistema 2026", font=font_small, fill=(0x94, 0xA3, 0xB8))
draw.text((80, 470), "Nivel: Avanzado", font=font_small, fill=(0xD4, 0xAF, 0x37))
draw.text((80, 600), "Manuel Francisco Machado\nLíder / Desarrollador", font=font_tiny, fill=(0x94, 0xA3, 0xB8))
img.save(str(portada_path))
print(f"Guardado portada: {portada_path}")

print("\nListo. Archivos generados en recursos/")
