"""
generar_presentacion_v2.py
Genera una presentación profesional tipo dashboard para ALBA.
Utiliza gráficos reales, tarjetas KPI, iconos y diseño oscuro+dorado premium.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from PIL import Image, ImageDraw, ImageFont
import comtypes.client

ROOT = Path(__file__).resolve().parent.parent
RESOURCES = ROOT / "recursos"
RESOURCES.mkdir(exist_ok=True)

# Palette ALBA premium
BG = RGBColor(0x05, 0x08, 0x13)
PANEL = RGBColor(0x0A, 0x0F, 0x1F)
GOLD = RGBColor(0xD4, 0xAF, 0x37)
GOLD_LIGHT = RGBColor(0xF0, 0xD7, 0x78)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SLATE = RGBColor(0x94, 0xA3, 0xB8)
SLATE_DARK = RGBColor(0x64, 0x70, 0x80)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
CYAN = RGBColor(0x06, 0xB6, 0xD4)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
PURPLE = RGBColor(0xA8, 0x55, 0xF7)
ROSE = RGBColor(0xF9, 0x73, 0x16)

ICONS = {
    "obs": "📊",
    "pred": "🔮",
    "match": "🎯",
    "emp": "🚀",
    "coach": "🧠",
    "sim": "⚗️",
    "data": "🗃️",
    "ai": "🤖",
    "impact": "💡",
    "repo": "🔗",
}


def set_slide_bg(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, int(slide_width_emu), int(slide_height_emu))
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    # send to back
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def add_rounded_panel(slide, left, top, width, height, fill=PANEL, line=GOLD, line_width=Pt(1), corner=0.12):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = line_width
    shape.adjustments[0] = corner
    return shape


def add_text(slide, text, left, top, width, height, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return box


def add_bullets(slide, items, left, top, width, height, size=16, color=SLATE, line_spacing=1.3):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(10)
        p.line_spacing = line_spacing
    return box


def add_footer(slide, text="ALBA — Datos al Ecosistema 2026 · Nivel Avanzado"):
    add_text(slide, text, 0.4, 7.05, 12.5, 0.3, size=10, color=GOLD, align=PP_ALIGN.RIGHT)


def add_kpi_card(slide, left, top, width, height, value, label, color=GOLD, subtext=""):
    panel = add_rounded_panel(slide, left, top, width, height, fill=PANEL, line=color, line_width=Pt(1.2))
    add_text(slide, value, left + 0.12, top + 0.12, width - 0.24, 0.55, size=32, color=color, bold=True)
    add_text(slide, label, left + 0.12, top + 0.62, width - 0.24, 0.35, size=12, color=SLATE, bold=True)
    if subtext:
        add_text(slide, subtext, left + 0.12, top + 0.95, width - 0.24, 0.25, size=10, color=SLATE_DARK)
    return panel


def add_module_card(slide, left, top, width, height, icon, title, desc, color):
    panel = add_rounded_panel(slide, left, top, width, height, fill=PANEL, line=color, line_width=Pt(1))
    add_text(slide, icon, left + 0.1, top + 0.08, 0.5, 0.45, size=28, align=PP_ALIGN.CENTER)
    add_text(slide, title, left + 0.1, top + 0.5, width - 0.2, 0.35, size=14, color=WHITE, bold=True)
    add_text(slide, desc, left + 0.1, top + 0.85, width - 0.2, height - 0.95, size=11, color=SLATE)
    # accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top + 0.04), Inches(0.08), Inches(height - 0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return panel


def add_colored_bar(slide, left, top, width, height, color):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def make_chart(slide, left, top, width, height, chart_type, chart_data, title=""):
    chart = slide.shapes.add_chart(
        chart_type,
        Inches(left), Inches(top), Inches(width), Inches(height),
        chart_data,
    ).chart
    chart.has_title = True if title else False
    if title:
        chart.chart_title.text_frame.text = title
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
        chart.chart_title.text_frame.paragraphs[0].font.color.rgb = WHITE
    chart.value_axis.visible = False
    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.category_axis.tick_labels.font.color.rgb = SLATE
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    return chart


# Slide dimensions
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide_width_emu = prs.slide_width
slide_height_emu = prs.slide_height
blank = prs.slide_layouts[6]

# ============================================================
# Slide 1: PORTADA
# ============================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s)
# Decorative golden diagonal line
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
line.fill.solid()
line.fill.fore_color.rgb = GOLD
line.line.fill.background()
# Title
add_text(s, "ALBA", 0.6, 1.5, 12, 1.2, size=90, color=GOLD, bold=True)
add_text(s, "Analítica Laboral Basada en IA", 0.6, 2.7, 12, 0.7, size=36, color=WHITE)
add_text(s, "Plataforma de inteligencia laboral para Colombia", 0.6, 3.4, 12, 0.5, size=18, color=SLATE)
# Info panel
add_rounded_panel(s, 0.6, 4.4, 5.5, 2.2, fill=PANEL, line=GOLD, line_width=Pt(1.5))
add_text(s, "Reto 5: Economía y Empleo", 0.8, 4.6, 5, 0.4, size=16, color=GOLD, bold=True)
add_text(s, "Concurso Datos al Ecosistema 2026 — IA para Colombia", 0.8, 5.0, 5, 0.4, size=14, color=WHITE)
add_text(s, "Nivel: Avanzado", 0.8, 5.45, 5, 0.35, size=14, color=SLATE)
add_text(s, "Manuel Francisco Machado\nLíder / Desarrollador", 0.8, 5.9, 5, 0.55, size=13, color=SLATE)
# Right side mini dashboard preview
add_rounded_panel(s, 7.2, 1.5, 5.5, 5.2, fill=PANEL, line=SLATE_DARK, line_width=Pt(1))
add_text(s, "En números", 7.45, 1.75, 5, 0.4, size=16, color=GOLD, bold=True)
kpis = [
    ("11", "Fuentes de datos"),
    ("44", "Tablas en Supabase"),
    ("744K", "Filas procesadas"),
    ("33", "Departamentos"),
]
for i, (v, l) in enumerate(kpis):
    y = 2.35 + i * 0.95
    add_text(s, v, 7.5, y, 1.5, 0.5, size=28, color=GOLD, bold=True)
    add_text(s, l, 9.0, y + 0.12, 3.5, 0.4, size=13, color=WHITE)
    s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(7.5), Inches(y + 0.6), Inches(12.3), Inches(y + 0.6)).line.color.rgb = SLATE_DARK
add_footer(s)

# ============================================================
# Slide 2: EL PROBLEMA (KPI cards + insight)
# ============================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_text(s, "El problema", 0.5, 0.35, 12, 0.8, size=42, color=WHITE, bold=True)
add_text(s, "Colombia tiene datos laborales, pero no una inteligencia que los una", 0.5, 1.05, 12, 0.4, size=16, color=SLATE)
# KPI row
kpis = [
    ("24M", "Ocupados", "GEIH 2025", GREEN),
    ("58%", "Informalidad", "GEIH 2025", ROSE),
    ("428K", "Programas SNIES", "Matriculados", CYAN),
    ("406", "Ocupaciones con salario", "DANE GEIH", GOLD),
]
for i, (v, l, sub, c) in enumerate(kpis):
    x = 0.5 + i * 3.15
    add_kpi_card(s, x, 1.65, 2.85, 1.35, v, l, c, sub)
# Insight panel
add_rounded_panel(s, 0.5, 3.25, 12.3, 3.6, fill=PANEL, line=GOLD, line_width=Pt(1))
add_text(s, "Desconexión estructural", 0.75, 3.45, 11.8, 0.4, size=20, color=GOLD, bold=True)
add_bullets(s, [
    "Ciudadanos no saben qué estudiar ni dónde hay trabajo.",
    "Emprendedores no saben qué negocio tiene potencial en su municipio.",
    "Empresas no encuentran talento con las habilidades que necesitan.",
    "Universidades desconocen si sus programas responden al mercado.",
    "El Gobierno diseña políticas sin visibilidad integral del territorio.",
], 0.75, 3.95, 11.8, 2.8, size=17, color=SLATE, line_spacing=1.4)
add_footer(s)

# ============================================================
# Slide 3: LA SOLUCIÓN — 6 módulos
# ============================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_text(s, "La solución", 0.5, 0.35, 12, 0.8, size=42, color=WHITE, bold=True)
add_text(s, "6 módulos que cubren todo el ciclo de decisión laboral", 0.5, 1.05, 12, 0.4, size=16, color=SLATE)
modules = [
    (ICONS["obs"], "Observatorio", "¿Qué está pasando en el mercado laboral?", BLUE),
    (ICONS["pred"], "Predicción IA", "¿Qué pasará a 5 y 10 años?", PURPLE),
    (ICONS["match"], "Match Inteligente", "¿Dónde encajo según mi perfil?", GOLD),
    (ICONS["emp"], "Emprende IA", "¿Qué negocio tiene potencial?", GREEN),
    (ICONS["coach"], "Coach IA", "¿Cómo me preparo para el empleo?", CYAN),
    (ICONS["sim"], "Simulación", "¿Qué pasa si cambio de carrera?", ROSE),
]
for i, (icon, title, desc, color) in enumerate(modules):
    col = i % 3
    row = i // 3
    x = 0.5 + col * 4.15
    y = 1.65 + row * 2.65
    add_module_card(s, x, y, 3.85, 2.35, icon, title, desc, color)
# Flow arrow bar at bottom
add_text(s, "Observar → Anticipar → Conectar → Emprender → Preparar → Simular", 0.5, 6.9, 12.3, 0.3, size=13, color=SLATE, align=PP_ALIGN.CENTER)
add_footer(s)

# ============================================================
# Slide 4: DATOS ABIERTOS — chart + table
# ============================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_text(s, "Datos abiertos utilizados", 0.5, 0.35, 12, 0.8, size=42, color=WHITE, bold=True)
add_text(s, "12 fuentes oficiales · 44 tablas · ~744.000 filas en Supabase", 0.5, 1.05, 12, 0.4, size=16, color=SLATE)
# Chart: rows per source
cd = ChartData()
cd.categories = ["GEIH", "SNIES", "ESCO", "RUES", "OLE", "SENA", "DNP", "PILA", "O*NET", "Saber Pro", "EMICRON", "World Bank"]
cd.add_series("Filas (miles)", [120, 85, 155, 26, 37, 17, 22, 0.652, 12, 1.4, 0.109, 0.128])
chart = make_chart(s, 0.5, 1.7, 8.0, 5.0, XL_CHART_TYPE.BAR_CLUSTERED, cd, "Volumen por fuente (miles de filas)")
# Style chart dark
plot = chart.plots[0]
series = chart.series[0]
series.format.fill.solid()
series.format.fill.fore_color.rgb = GOLD
chart.category_axis.tick_labels.font.size = Pt(10)
chart.category_axis.tick_labels.font.color.rgb = SLATE
chart.value_axis.tick_labels.font.color.rgb = SLATE
# Source detail panel
add_rounded_panel(s, 8.8, 1.7, 4.0, 5.0, fill=PANEL, line=GOLD, line_width=Pt(1))
add_text(s, "Fuentes principales", 9.0, 1.9, 3.6, 0.35, size=14, color=GOLD, bold=True)
details = [
    ("DANE GEIH", "Empleo, salarios, informalidad"),
    ("MEN SNIES + OLE", "Oferta educativa e ingresos"),
    ("SENA SPE/APE", "Demanda laboral y cursos"),
    ("MinTrabajo PILA", "Empleo formal por sector"),
    ("Confecámaras RUES", "Empresas nuevas"),
    ("ESCO + O*NET", "Habilidades y ocupaciones"),
    ("World Bank", "Indicadores macro"),
]
for i, (src, desc) in enumerate(details):
    y = 2.4 + i * 0.58
    add_text(s, src, 9.0, y, 3.6, 0.28, size=11, color=WHITE, bold=True)
    add_text(s, desc, 9.0, y + 0.22, 3.6, 0.28, size=10, color=SLATE)
add_footer(s)

# ============================================================
# Slide 5: ARQUITECTURA
# ============================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_text(s, "Arquitectura tecnológica", 0.5, 0.35, 12, 0.8, size=42, color=WHITE, bold=True)
add_text(s, "Full-stack conectado a datos oficiales + modelos de IA", 0.5, 1.05, 12, 0.4, size=16, color=SLATE)
# Architecture stack
layers = [
    ("Frontend", "React 18 · TypeScript · Vite · Tailwind", BLUE),
    ("Backend", "FastAPI · Uvicorn · 6 routers · 40+ endpoints", PURPLE),
    ("Base de datos", "Supabase PostgreSQL + pgvector · 44 tablas", GOLD),
    ("IA / Forecasting", "Gemini 2.5 Flash-Lite · Gemini Live · Chronos T5 · Gemma 300", GREEN),
    ("Fuentes", "DANE · MEN · SENA · MinTrabajo · DNP · RUES · ESCO · O*NET · World Bank", SLATE),
]
for i, (title, desc, color) in enumerate(layers):
    y = 1.65 + i * 1.05
    panel = add_rounded_panel(s, 0.6, y, 12.1, 0.85, fill=PANEL, line=color, line_width=Pt(1.5))
    add_text(s, title, 0.9, y + 0.18, 2.8, 0.5, size=16, color=color, bold=True)
    add_text(s, desc, 3.9, y + 0.24, 8.5, 0.5, size=13, color=WHITE)
# Small connector lines between layers
for i in range(len(layers) - 1):
    y1 = 1.65 + i * 1.05 + 0.85
    line = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(6.4), Inches(y1), Inches(6.4), Inches(y1 + 0.2))
    line.line.color.rgb = SLATE_DARK
    line.line.width = Pt(2)
add_footer(s)

# ============================================================
# Slide 6: RESULTADOS — dashboard con gráfico
# ============================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_text(s, "Resultados y demo", 0.5, 0.35, 12, 0.8, size=42, color=WHITE, bold=True)
add_text(s, "De datos dispersos a decisiones accionables", 0.5, 1.05, 12, 0.4, size=16, color=SLATE)
# Mini KPI row
mini_kpis = [
    ("6", "Módulos", BLUE),
    ("40+", "Endpoints", PURPLE),
    ("406", "Salarios reales", GOLD),
    ("21", "Profesiones proyectadas", GREEN),
]
for i, (v, l, c) in enumerate(mini_kpis):
    x = 0.5 + i * 3.15
    add_kpi_card(s, x, 1.65, 2.85, 1.1, v, l, c)
# Chart: projected sectors growth
cd2 = ChartData()
cd2.categories = ["Tecnología", "Salud", "Energía", "Construcción", "Comercio", "Manufactura", "Agricultura"]
cd2.add_series("Crecimiento proyectado 2035 (%)", [38, 31, 25, 18, 14, 10, 7])
chart2 = make_chart(s, 0.5, 3.0, 7.5, 3.8, XL_CHART_TYPE.COLUMN_CLUSTERED, cd2, "Sectores con mayor proyección")
chart2.series[0].format.fill.solid()
chart2.series[0].format.fill.fore_color.rgb = GOLD
chart2.category_axis.tick_labels.font.color.rgb = SLATE
chart2.value_axis.tick_labels.font.color.rgb = SLATE
# Highlight bullets
add_rounded_panel(s, 8.3, 3.0, 4.5, 3.8, fill=PANEL, line=GOLD, line_width=Pt(1))
add_text(s, "Hallazgos clave", 8.5, 3.2, 4.1, 0.35, size=15, color=GOLD, bold=True)
add_bullets(s, [
    "Predicciones a 5 y 10 años para sectores, profesiones y salarios.",
    "Salarios reales de 406 ocupaciones del DANE GEIH.",
    "Matching híbrido: 50% ESCO + 50% LLM.",
    "Mapa territorial con empleo por departamento y sector.",
    "Botón 'Analizar con IA' en cada gráfica.",
], 8.5, 3.65, 4.1, 3.0, size=13, color=SLATE, line_spacing=1.35)
add_footer(s)

# ============================================================
# Slide 7: IMPACTO
# ============================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_text(s, "Impacto esperado", 0.5, 0.35, 12, 0.8, size=42, color=WHITE, bold=True)
add_text(s, "Beneficiarios directos del uso de datos abiertos + IA", 0.5, 1.05, 12, 0.4, size=16, color=SLATE)
impact = [
    ("👤", "Ciudadanos", "Saben qué estudiar, dónde hay trabajo y cuánto ganar.", BLUE),
    ("🏢", "Empresas", "Encuentran talento y entienden la demanda de su sector.", PURPLE),
    ("🎓", "Universidades", "Alinean programas con la demanda laboral real.", GOLD),
    ("🚀", "Emprendedores", "Identifican negocios con potencial en su municipio.", GREEN),
    ("🏛️", "Gobierno", "Diseña políticas con visibilidad territorial integral.", CYAN),
    ("📈", "Escalabilidad", "Arquitectura lista para Vercel + Railway + Supabase.", ROSE),
]
for i, (icon, title, desc, color) in enumerate(impact):
    col = i % 3
    row = i // 3
    x = 0.5 + col * 4.15
    y = 1.65 + row * 2.65
    add_module_card(s, x, y, 3.85, 2.35, icon, title, desc, color)
add_footer(s)

# ============================================================
# Slide 8: CIERRE
# ============================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_text(s, "¿Por qué ALBA?", 0.5, 0.35, 12, 0.8, size=42, color=WHITE, bold=True)
add_rounded_panel(s, 0.5, 1.4, 12.3, 4.8, fill=PANEL, line=GOLD, line_width=Pt(1.5))
add_text(s, "Valor diferencial", 0.75, 1.65, 11.8, 0.45, size=22, color=GOLD, bold=True)
add_bullets(s, [
    "Integra por primera vez 11 fuentes de datos laborales y educativas en una sola plataforma.",
    "Usa IA explicable: cada recomendación muestra los datos reales que la soportan.",
    "Análisis territorial real: empleo, salarios y sectores por cada departamento.",
    "Forecasting con Chronos T5 y salarios reales del DANE, no estimaciones genéricas.",
    "Código abierto, documentado y listo para desplegar en la nube.",
], 0.75, 2.25, 11.8, 3.0, size=17, color=SLATE, line_spacing=1.4)
# Links panel
add_rounded_panel(s, 0.5, 6.4, 12.3, 0.8, fill=PANEL, line=SLATE_DARK, line_width=Pt(1))
add_text(s, "🌐 Demo · 📂 Repositorio · 📊 Presentación · 📄 Documentación", 0.75, 6.58, 11.8, 0.4, size=14, color=WHITE, align=PP_ALIGN.CENTER)
add_footer(s, "ALBA — Analítica Laboral Basada en IA · Datos al Ecosistema 2026")

# Save PPTX
pptx_path = RESOURCES / "presentacion.pptx"
prs.save(str(pptx_path))
print(f"PPTX guardado: {pptx_path}")

# Export to PDF
pdf_path = RESOURCES / "presentacion.pdf"
try:
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1
    deck = powerpoint.Presentations.Open(str(pptx_path), WithWindow=False)
    deck.SaveAs(str(pdf_path), 32)
    deck.Close()
    powerpoint.Quit()
    print(f"PDF guardado: {pdf_path}")
except Exception as e:
    print(f"Error exportando PDF: {e}")

# Generate cover image
portada_path = RESOURCES / "portada.png"
img = Image.new("RGB", (1920, 1080), color=(0x05, 0x08, 0x13))
draw = ImageDraw.Draw(img)
try:
    font_dir = Path("C:/Windows/Fonts")
    f_title = ImageFont.truetype(str(font_dir / "Calibri.ttf"), 160)
    f_sub = ImageFont.truetype(str(font_dir / "Calibri.ttf"), 52)
    f_body = ImageFont.truetype(str(font_dir / "Calibri.ttf"), 34)
    f_small = ImageFont.truetype(str(font_dir / "Calibri.ttf"), 26)
except Exception:
    f_title = f_sub = f_body = f_small = ImageFont.load_default()

# Left accent bar
draw.rectangle([0, 0, 20, 1080], fill=(0xD4, 0xAF, 0x37))
draw.rectangle([40, 280, 260, 292], fill=(0xD4, 0xAF, 0x37))
draw.text((80, 140), "ALBA", font=f_title, fill=(0xD4, 0xAF, 0x37))
draw.text((80, 340), "Analítica Laboral Basada en IA", font=f_sub, fill=(0xFF, 0xFF, 0xFF))
draw.text((80, 430), "Plataforma de inteligencia laboral para Colombia", font=f_body, fill=(0x94, 0xA3, 0xB8))
draw.text((80, 520), "Reto 5: Economía y Empleo — Concurso Datos al Ecosistema 2026", font=f_small, fill=(0x94, 0xA3, 0xB8))
draw.text((80, 580), "Nivel: Avanzado", font=f_small, fill=(0xD4, 0xAF, 0x37))
draw.text((80, 720), "Manuel Francisco Machado\nLíder / Desarrollador", font=f_small, fill=(0x94, 0xA3, 0xB8))

# Right KPI dashboard preview
panel = [(1260, 150), (1860, 150), (1860, 750), (1260, 750)]
for i in range(len(panel)):
    x1, y1 = panel[i]
    x2, y2 = panel[(i + 1) % len(panel)]
    draw.line([(x1, y1), (x2, y2)], fill=(0xD4, 0xAF, 0x37), width=2)
draw.text((1300, 200), "En números", font=f_body, fill=(0xD4, 0xAF, 0x37))
numbers = [("11", "Fuentes"), ("44", "Tablas"), ("744K", "Filas"), ("33", "Deptos")]
for i, (n, l) in enumerate(numbers):
    y = 290 + i * 110
    draw.text((1300, y), n, font=f_sub, fill=(0xD4, 0xAF, 0x37))
    draw.text((1450, y + 12), l, font=f_body, fill=(0xFF, 0xFF, 0xFF))
    draw.line([(1300, y + 70), (1820, y + 70)], fill=(0x64, 0x70, 0x80), width=1)

img.save(str(portada_path))
print(f"Portada guardada: {portada_path}")
print("\nPresentacion rediseniada exitosamente.")
