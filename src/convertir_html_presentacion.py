"""
Convierte alba_slides.html a:
  - presentacion.pptx (imágenes de cada diapositiva)
  - presentacion.pdf
  - portada.png
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import comtypes.client

ROOT = Path(__file__).resolve().parent.parent
RESOURCES = ROOT / "recursos"
HTML_FILE = RESOURCES / "alba_slides.html"
PREVIEWS = RESOURCES / "slide_previews_html"
PREVIEWS.mkdir(exist_ok=True)

# 1. Tomar screenshots con Playwright
from playwright.sync_api import sync_playwright

with sync_playwright() as p:
    browser = p.chromium.launch()
    page = browser.new_page(viewport={"width": 1920, "height": 1080})
    page.goto(f"file:///{HTML_FILE.resolve().as_posix()}")
    page.wait_for_timeout(1000)

    # detect total slides from JS
    total = page.evaluate("document.querySelectorAll('.slide').length")
    print(f"Total diapositivas detectadas: {total}")

    for i in range(total):
        # navigate to slide i
        page.evaluate(f"goTo({i})")
        page.wait_for_timeout(600)
        path = PREVIEWS / f"slide_{i+1:02d}.png"
        page.screenshot(path=str(path), full_page=False)
        print(f"Screenshot {i+1}/{total}")

    browser.close()

# 2. Crear PPTX a partir de screenshots
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

for i in range(1, total + 1):
    img_path = PREVIEWS / f"slide_{i:02d}.png"
    slide = prs.slides.add_slide(blank)
    # Add image covering full slide
    slide.shapes.add_picture(str(img_path), Inches(0), Inches(0), width=Inches(13.333), height=Inches(7.5))

pptx_path = RESOURCES / "presentacion.pptx"
prs.save(str(pptx_path))
print(f"PPTX guardado: {pptx_path}")

# 3. Exportar a PDF con PowerPoint
pdf_path = RESOURCES / "presentacion.pdf"
try:
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1
    deck = powerpoint.Presentations.Open(str(pptx_path), WithWindow=False)
    deck.SaveAs(str(pdf_path), 32)
    deck.Close()
    powerpoint.Quit()
    print(f"PDF guardado: {pdf_path}")
except Exception as e:
    print(f"Error exportando PDF: {e}")

# 4. Portada desde primer screenshot
portada_path = RESOURCES / "portada.png"
img = Image.open(str(PREVIEWS / "slide_01.png"))
img.save(str(portada_path))
print(f"Portada guardada: {portada_path}")

print("\nPresentacion web convertida a PPTX, PDF y PNG exitosamente.")
